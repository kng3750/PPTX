"""
PPTX 템플릿에 문서 내용을 삽입하는 웹 애플리케이션
- 템플릿의 배경색, 머릿말(타이틀)은 유지
- 본문 영역에만 문서 내용 삽입
"""

import copy
import io
import os
from pathlib import Path

from flask import Flask, render_template, request, send_file
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from docx import Document as DocxDocument

app = Flask(__name__)
app.config["MAX_CONTENT_LENGTH"] = 16 * 1024 * 1024  # 16MB
UPLOAD_FOLDER = Path("uploads")
UPLOAD_FOLDER.mkdir(exist_ok=True)

# 머릿말/타이틀로 유지할 placeholder 타입
HEADER_TYPES = {
    PP_PLACEHOLDER.TITLE,
    PP_PLACEHOLDER.CENTER_TITLE,
    PP_PLACEHOLDER.SUBTITLE,
    PP_PLACEHOLDER.VERTICAL_TITLE,
}

# 본문으로 처리할 placeholder 타입
BODY_TYPES = {
    PP_PLACEHOLDER.BODY,
    PP_PLACEHOLDER.VERTICAL_BODY,
    PP_PLACEHOLDER.OBJECT,  # 일부 템플릿에서 본문이 OBJECT로 정의됨
}


def extract_text_from_document(file_path: Path, filename: str) -> str:
    """문서 파일에서 텍스트 추출"""
    ext = Path(filename).suffix.lower()

    if ext == ".txt":
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            return f.read()

    if ext in (".docx", ".doc"):
        doc = DocxDocument(file_path)
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text)
        return "\n".join(paragraphs)

    raise ValueError(f"지원하지 않는 파일 형식입니다: {ext}")


def _find_body_placeholder(slide):
    """슬라이드에서 본문 placeholder 찾기"""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        try:
            if shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                if ph_type in HEADER_TYPES:
                    continue
                if ph_type in BODY_TYPES:
                    return shape
        except (ValueError, AttributeError):
            continue
    # BODY 타입이 없으면 TITLE이 아닌 텍스트 shape 중 가장 큰 것
    candidates = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        try:
            if shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                if ph_type in HEADER_TYPES:
                    continue
            area = shape.width * shape.height
            candidates.append((area, shape))
        except (ValueError, AttributeError):
            continue
    if candidates:
        candidates.sort(key=lambda x: -x[0])
        return candidates[0][1]
    return None


def _duplicate_slide(prs, slide_index: int):
    """슬라이드 복제 (배경·머릿말 포함)"""
    template = prs.slides[slide_index]
    try:
        blank_layout = prs.slide_layouts[6]
    except IndexError:
        blank_layout = prs.slide_layouts[len(prs.slide_layouts) - 1]
    new_slide = prs.slides.add_slide(blank_layout)
    for shp in template.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(newel, "p:extLst")
    for _, value in template.part.rels.items():
        if "notesSlide" not in value.reltype:
            try:
                new_slide.part.rels.add_relationship(
                    value.reltype, value._target, value.rId
                )
            except Exception:
                pass
    return new_slide


def _fill_body_with_text(body_placeholder, lines: list, center_align: bool = True):
    """본문 placeholder에 텍스트 삽입 (중앙 정렬 옵션)"""
    tf = body_placeholder.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        if center_align:
            p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER


def insert_document_into_pptx(
    template_path: Path, doc_text: str, num_slides: int = 1
) -> io.BytesIO:
    """
    PPTX 템플릿의 본문 영역에만 문서 내용을 삽입.
    배경, 머릿말(타이틀)은 변경하지 않음.
    num_slides: 생성할 슬라이드 수 (선택한 페이지 수)
    """
    prs = Presentation(str(template_path))
    slide = prs.slides[0]
    body_placeholder = _find_body_placeholder(slide)

    if body_placeholder is None:
        raise ValueError(
            "템플릿에서 본문 영역을 찾을 수 없습니다. "
            "본문 placeholder(텍스트 상자)가 있는 슬라이드를 사용해 주세요."
        )

    lines = [ln.strip() for ln in doc_text.strip().split("\n") if ln.strip()]
    if not lines:
        lines = [doc_text.strip()] if doc_text.strip() else [""]

    num_slides = max(1, min(40, int(num_slides)))
    # 슬라이드 수만큼 복제 (이미 1장 있으므로 num_slides-1장 추가)
    for _ in range(num_slides - 1):
        _duplicate_slide(prs, 0)

    # 문서 내용을 슬라이드 수만큼 균등 분할
    n = len(lines)
    chunk_size = (n + num_slides - 1) // num_slides if num_slides else n
    chunks = []
    for i in range(num_slides):
        start = i * chunk_size
        end = min((i + 1) * chunk_size, n)
        chunks.append(lines[start:end] if start < n else [""])

    for idx, slide in enumerate(prs.slides):
        body = _find_body_placeholder(slide)
        if body and idx < len(chunks):
            _fill_body_with_text(body, chunks[idx], center_align=True)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


@app.route("/")
def index():
    return render_template("index.html")


@app.route("/process", methods=["POST"])
def process():
    if "template" not in request.files or "document" not in request.files:
        return {"error": "템플릿 파일과 문서 파일을 모두 업로드해 주세요."}, 400

    template_file = request.files["template"]
    document_file = request.files["document"]

    if template_file.filename == "" or document_file.filename == "":
        return {"error": "템플릿 파일과 문서 파일을 모두 선택해 주세요."}, 400

    if not template_file.filename.lower().endswith(".pptx"):
        return {"error": "템플릿은 .pptx 파일이어야 합니다."}, 400

    doc_ext = Path(document_file.filename).suffix.lower()
    if doc_ext not in (".txt", ".docx", ".doc"):
        return {"error": "문서는 .txt, .docx, .doc 파일만 지원합니다."}, 400

    try:
        # 임시 저장
        template_path = UPLOAD_FOLDER / f"template_{os.urandom(8).hex()}.pptx"
        doc_path = UPLOAD_FOLDER / f"doc_{os.urandom(8).hex()}{doc_ext}"
        template_file.save(template_path)
        document_file.save(doc_path)

        try:
            doc_text = extract_text_from_document(doc_path, document_file.filename)
            if not doc_text.strip():
                return {"error": "문서에서 추출할 텍스트가 없습니다."}, 400

            # 선택한 페이지 수 (체크된 개수 = 생성할 슬라이드 수)
            selected_pages = request.form.getlist("pages")
            num_slides = len(selected_pages) if selected_pages else 1

            output = insert_document_into_pptx(
                template_path, doc_text, num_slides=num_slides
            )
            return send_file(
                output,
                as_attachment=True,
                download_name="result.pptx",
                mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        finally:
            template_path.unlink(missing_ok=True)
            doc_path.unlink(missing_ok=True)
    except ValueError as e:
        return {"error": str(e)}, 400
    except Exception as e:
        return {"error": f"처리 중 오류가 발생했습니다: {str(e)}"}, 500


if __name__ == "__main__":
    app.run(debug=True, port=5000)
